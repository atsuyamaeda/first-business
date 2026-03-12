#!/usr/bin/env python3
"""抹茶スタートアップ ピッチデッキ生成スクリプト v2 (World-Class Design)
Design language: Figma / Linear / Vercel aesthetic
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

from reportlab.lib.units import cm, mm
from reportlab.lib import colors as rl_colors
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.cidfonts import UnicodeCIDFont
from reportlab.pdfgen import canvas as pdf_canvas

# ── カラーパレット ──────────────────────────────────────────
# Dark slides (cover, ask, closing)
DARK1   = RGBColor(0x09, 0x0D, 0x18)   # near-black navy
DARK2   = RGBColor(0x14, 0x24, 0x3B)   # gradient end (lighter)
GREEN   = RGBColor(0x1D, 0xB8, 0x6A)   # vibrant matcha green
GDARK   = RGBColor(0x0E, 0x7A, 0x47)   # dark green for text on light
GOLD    = RGBColor(0xF0, 0xB4, 0x29)   # warm amber gold
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
DIM     = RGBColor(0x7A, 0x95, 0xB5)   # muted text on dark
DECO_G  = RGBColor(0x0B, 0x23, 0x15)   # decorative circle dark green
DECO_B  = RGBColor(0x07, 0x14, 0x2B)   # decorative circle dark blue

# Light slides (content, bigstat, table, threeCol, threeStep)
LBKG    = RGBColor(0xFF, 0xFF, 0xFF)   # pure white
DTXT    = RGBColor(0x0F, 0x17, 0x2A)   # near-black text
MTXT    = RGBColor(0x47, 0x55, 0x69)   # medium gray
BRDR    = RGBColor(0xE1, 0xE8, 0xF0)   # border / separator
CARD_G  = RGBColor(0xF0, 0xFD, 0xF4)   # light green card bg
CARD_B  = RGBColor(0xEF, 0xF6, 0xFF)   # light blue card bg
CARD_A  = RGBColor(0xFF, 0xFB, 0xEB)   # light amber card bg
BLUE_H  = RGBColor(0x21, 0x63, 0xD0)   # blue accent

# reportlab counterparts
RL_DARK1  = rl_colors.HexColor('#090D18')
RL_GREEN  = rl_colors.HexColor('#1DB86A')
RL_GDARK  = rl_colors.HexColor('#0E7A47')
RL_GOLD   = rl_colors.HexColor('#F0B429')
RL_WHITE  = rl_colors.white
RL_DIM    = rl_colors.HexColor('#7A95B5')
RL_DTXT   = rl_colors.HexColor('#0F172A')
RL_MTXT   = rl_colors.HexColor('#475569')
RL_BRDR   = rl_colors.HexColor('#E1E8F0')
RL_CARD_G = rl_colors.HexColor('#F0FDF4')
RL_CARD_B = rl_colors.HexColor('#EFF6FF')
RL_CARD_A = rl_colors.HexColor('#FFFBEB')
RL_BLUE_H = rl_colors.HexColor('#2163D0')
RL_LBKG   = rl_colors.white

# ── スライドサイズ (16:9) ────────────────────────────────────
SW   = Inches(13.33)
SH   = Inches(7.5)
SWcm = 33.87
SHcm = 19.05

# ── reportlab ページサイズ ───────────────────────────────────
PAGE_W = 33.87 * cm
PAGE_H = 19.05 * cm

# ── フォント ─────────────────────────────────────────────────
pdfmetrics.registerFont(UnicodeCIDFont('HeiseiKakuGo-W5'))
pdfmetrics.registerFont(UnicodeCIDFont('HeiseiMin-W3'))
JP_FONT  = 'HeiseiKakuGo-W5'
JP_SERIF = 'HeiseiMin-W3'

# ─────────────────────────────────────────────────────────────
# スライドデータ定義
# ─────────────────────────────────────────────────────────────
SLIDES = [
  # 1. カバー
  {
    "type": "cover",
    "title": "[COMPANY]",
    "subtitle": "The Nespresso of Matcha",
    "body": [
      "シードラウンド  ·  $3,000,000  ·  機密",
      "2026年3月",
      "",
      "「800年前、禅僧が『養生の仙薬』と記した。",
      " 2024年、科学がそれを証明した。",
      " 3,000万台のNespressoマシンで、30秒で。」",
    ],
    "note": "カバースライド"
  },
  # 2. 投資テーゼ
  {
    "type": "content",
    "tag": "INVESTMENT THESIS",
    "title": "1文で伝える投資機会",
    "lead": "[COMPANY]は、北米3,000万台以上のNespressoオーナー向けに\n世界初の日本産セレモニアル抹茶サブスクリプションを構築する。",
    "body": [
      "① カテゴリの転換、トレンドではない — 「抹茶 vs コーヒー」検索数：+300〜500%（2020年〜）",
      "② 市場は壊れており、私たちが修復する — $6の粗悪品 vs $30+の高嶺の花。準備の手間が初回購入者の70%を30日以内に失わせる",
      "③ 成功モデルが実証済み — PerfectTed（英国）はNespresso対応抹茶ポッドで4年間CAGR 532%、£140Mバリュエーションを達成。彼らは英国のみ。北米は競合ゼロ",
    ],
    "note": "PerfectTedの比較を早期に提示することで『前例はあるか？』という質問に先手を打つ"
  },
  # 3. 問題① コーヒー依存
  {
    "type": "bigstat",
    "tag": "PROBLEM 1/3",
    "title": "米国人の68%がコーヒー不安・クラッシュを経験",
    "stats": [
      ("68%", "コーヒー起因の不安を定期的に経験"),
      ("45%", "コーヒーの代替を積極探索中"),
      ("39%", "Z世代がコーヒーを削減・断絶"),
    ],
    "body": [
      "消費者はカフェインを捨てない — より良いカフェインを求めている",
      "Z世代：TikTokで健康意識を持って育った初の世代",
      "コーヒーから抹茶へ「移行」するのではなく、コーヒーを「スキップ」している",
      "TikTok #matcha：約480億回再生 — 世界最大の食品・飲料カテゴリ",
    ],
  },
  # 4. 問題② 抹茶市場の崩壊
  {
    "type": "content",
    "tag": "PROBLEM 2/3",
    "title": "抹茶は答え — しかし市場は構造的に崩壊している",
    "lead": "科学的に証明された「穏やかな覚醒」",
    "body": [
      "L-テアニン（36〜50mg/杯）：40分以内にα波増加 → 「穏やかな集中」（Mason, 2001）",
      "EGCG：通常緑茶ティーバッグの137倍の抗酸化物質（Weiss & Anderton, 2003）",
      "2024年 PLOS ONE RCT（99名・12ヶ月）：認知機能の有意改善（p=0.028）",
      "",
      "【市場の3大崩壊】",
      "① 価格のミッシングミドル：$6のゴミ vs $30+の高嶺の花（$15〜20の品質デイリー粉末が存在しない）",
      "② 準備の壁：茶筅・特定温度・茶こし・抹茶碗が必要 → 初回購入者の70%が30日で離脱",
      "③ 品質の信頼危機：『セレモニアルグレード』は無規制のマーケティング用語",
    ],
  },
  # 5. 問題③ 準備の壁
  {
    "type": "bigstat",
    "tag": "PROBLEM 3/3",
    "title": "準備の複雑さが市場拡大の最大障壁",
    "stats": [
      ("70%", "初回購入者が30日以内に再購入しない"),
      ("7.3%", "Amazonネガティブレビューでダマ問題を言及"),
      ("37.4%", "価格不満（$6ゴミ vs $30+高嶺）"),
    ],
    "body": [
      "最も多い消費者の叫び：『自宅でカフェクオリティの抹茶を、手間なく飲みたい』",
      "現在の代替手段：電動ミルクフォーマー、シェイカーボトル、ブレンダー — すべて妥協の産物",
      "消費者は解決策にお金を払う意思がある（Pique: $2.43/杯でも売れている）",
    ],
  },
  # 6. ソリューション
  {
    "type": "content",
    "tag": "SOLUTION",
    "title": "日本産セレモニアル抹茶 × Nespresso対応ポッド",
    "lead": "ボタンを押すだけ → 30秒で完璧な抹茶",
    "body": [
      "✓  既存のNespressoマシンで動作（新規ハードウェア不要）",
      "✓  ダマなし・道具不要・完璧な温度管理",
      "✓  日本の農家から直接調達した認定セレモニアルグレード",
      "✓  窒素封入個包装 → 1杯分ずつ鮮度保証",
      "✓  QRコードで農園・収穫日・茶師の認定を確認可能",
      "",
      "Nespressoが1億人のコーヒー愛好家に革命を起こしたように、",
      "[COMPANY]は抹茶を30秒の日常習慣へ変革する",
    ],
  },
  # 7. 仕組み（3ステップ）
  {
    "type": "threeStep",
    "tag": "HOW IT WORKS",
    "title": "3ステップで完成",
    "steps": [
      ("STEP 1\nサブスク申込", "毎月ポッドをお届け\n配送料無料・いつでも一時停止可"),
      ("STEP 2\nポッドを挿入", "既存のNespresso\nまたは互換機に装着"),
      ("STEP 3\n30秒で完成", "完璧な温度・濃度の\n本物の日本産抹茶"),
    ],
  },
  # 8. プロダクトライン
  {
    "type": "table",
    "tag": "PRODUCT LINE",
    "title": "プロダクトラインナップ",
    "thead": ["SKU", "内容", "価格", "ポジション"],
    "rows": [
      ["CEREMONIAL POD", "宇治・京都産 セレモニアル  10ポッド", "$29", "最高品質・シングルオリジン"],
      ["DAILY POD",      "プレミアムグレード  10ポッド",         "$22", "毎日使いのデイリードライバー"],
      ["POWDER+",        "産地明記 30g 粉末缶",                  "$18", "ミッシングミドル埋め"],
      ["FLAGSHIP BUNDLE","30ポッド + 30g 粉末  月次サブスク",    "$65/月", "主力・最高 LTV"],
    ],
    "note": "ポッドが主力。粉末は入口商品として機能し、ポッドへの段階的誘導を設計"
  },
  # 9. 日本農家バリューチェーン
  {
    "type": "content",
    "tag": "SUPPLY CHAIN",
    "title": "日本農家直送 → ポッドへ",
    "lead": "サプライチェーンが私たちの最大の競争優位",
    "body": [
      "【調達先（目標）】",
      "・宇治（京都）：800年の権威。最上位セレモニアルグレード",
      "・西尾（愛知）：全国生産量1位（20〜30%シェア）。抹茶特化産地",
      "・鹿児島：急成長の革新産地。鹿児島産 2024年荒茶生産量全国1位",
      "・八女（福岡）：全国茶品評会玉露部門 18年連続産地賞",
      "",
      "【製品品質】",
      "・日本国内OEM製造 → 日本産の真正性を維持",
      "・認定茶師（茶師）がバッチごとに品質確認",
      "・ISO 20715:2023 準拠（2023年制定 国際抹茶標準）",
      "・QRコード → ブロックチェーントレーサビリティ（農園〜ポッド）",
    ],
  },
  # 10. 市場規模
  {
    "type": "bigstat",
    "tag": "MARKET SIZE",
    "title": "市場規模：TAM / SAM / SOM",
    "stats": [
      ("$4.3B", "TAM 世界抹茶市場（2024年）→ $8.9B（2034年）CAGR 7.8%"),
      ("$1.4B", "SAM 北米＋英国プレミアム抹茶市場  CAGR 11%（最速成長地域）"),
      ("$280M", "SOM プレミアム抹茶サブスク（ポッド＋品質粉末）"),
    ],
    "body": [
      "Year 3 目標：$14M ARR = SOM の 5% 獲得",
      "Nespresso対応カプセル市場：$6B+（抹茶ポッドはほぼ空白）",
    ],
  },
  # 11. なぜ今か
  {
    "type": "content",
    "tag": "WHY NOW",
    "title": "抹茶モーメント：ティッピングポイントを超えた",
    "lead": "これはトレンドではなく、カテゴリシフトである",
    "body": [
      "「matcha」Google検索：+292〜520%（2020年〜2025年）",
      "「matcha vs coffee」検索：+300〜500%（全カテゴリ最速）",
      "TikTok #matcha：約480億回再生",
      "米国小売抹茶販売：3年間で86%増（NielsenIQ）",
      "McDonald's オーストラリア：2025年10月 全国的に抹茶ライン投入",
      "Starbucks 2025：抹茶製品をリニューアル（メインストリーム統合）",
      "米国レストランの5.71%（213,996店）が抹茶メニュー提供（前年比+30%）",
      "",
      "→ インフラを持つ者が勝つ。今が参入の最後のウィンドウ",
    ],
  },
  # 12. 供給危機
  {
    "type": "content",
    "tag": "OUR MOAT",
    "title": "供給危機 = 私たちの競争優位",
    "lead": "京都の碾茶価格：1年で +265%（2024→2025）",
    "body": [
      "一保堂（創業1717年）：2024年10月 抹茶全販売を一時停止",
      "丸久小山園：新規卸売顧客の受付を停止",
      "米国卸売リードタイム：1〜2ヶ月 → 6ヶ月に延長",
      "「抹茶転売ヤー」が二次市場で初めて出現",
      "",
      "【なぜこの危機は長期化するか】",
      "・日本の茶農家人口：2010年以降 40%減少（高齢化・後継者不足）",
      "・新茶畑の成熟：5年かかる → 供給は需要に追いつけない",
      "・茶農家の平均年齢：67歳以上（構造的問題）",
      "",
      "【結論】今すぐサプライチェーン関係を築いた者が、次の10年を制する",
    ],
  },
  # 13. ターゲット顧客
  {
    "type": "content",
    "tag": "TARGET CUSTOMER",
    "title": "ターゲット顧客",
    "lead": "プライマリ：アッパーミドル・25〜40歳・都市部プロフェッショナル",
    "body": [
      "【プライマリ層の特徴】",
      "・世帯年収 $80K+。Nespressoマシンを既に所有（米国だけで3,000万台以上）",
      "・健康志向だが時間がない。コーヒーに$8〜15/日を既に支出",
      "・新規ハードウェアゼロ。ゼロの学習コスト",
      "",
      "【セカンダリ層の特徴】",
      "・健康意識の高いZ世代（25〜35歳）コーヒー離れ世代",
      "・Z世代の62%が抹茶を「集中力・認知機能向上」と関連付け",
      "・TikTok/Instagramで抹茶を発見し、準備方法に圧倒されている",
      "・41%がソーシャルアプリを主要検索エンジンとして使用",
    ],
  },
  # 14. 3消費者セグメント
  {
    "type": "threeCol",
    "tag": "CONSUMER SEGMENTS",
    "title": "3つの消費者セグメント、1つのソリューション",
    "cols": [
      ("DAILY DRINKERS\n市場の60%",
       "自宅でバリスタ品質を\n手間なく飲みたい\n\n→ ポッドが完全解決"),
      ("NEWCOMERS\n市場の25%",
       "TikTokで抹茶を発見\n準備方法に圧倒されている\n\n→ ポッドは完璧な入口"),
      ("PURISTS\n市場の15%",
       "日本産を既に購入済み\nプレミアムな便利さを求める\n\n→ 農園ストーリーで満足"),
    ],
  },
  # 15. 未充足ニーズ
  {
    "type": "table",
    "tag": "UNMET NEEDS",
    "title": "解決する5つの未充足ニーズ",
    "thead": ["#", "消費者の課題", "私たちの解決策"],
    "rows": [
      ["①", "準備が複雑すぎる（道具・温度・技術）", "ポッド = 30秒・道具不要"],
      ["②", "$6品質ゴミ vs $30+高嶺の花（ミッシングミドル）", "POWDER+ $18/30g で埋める"],
      ["③", "「セレモニアルグレード」は無規制で信頼できない", "ISO認証＋農園QRトレーサビリティ"],
      ["④", "開封後4〜8週間で鮮度劣化", "窒素封入個包装 = 1杯分ずつ新鮮"],
      ["⑤", "サブスクが解約できない罠", "いつでも一時停止・解約・スキップ可"],
    ],
  },
  # 16. PerfectTedの証明
  {
    "type": "content",
    "tag": "VALIDATION",
    "title": "PerfectTed：このモデルが機能することの証明",
    "lead": "英国で全く同じモデルを4年間で CAGR 532% 達成",
    "body": [
      "2021年創業。自己資金：£250K + 家族融資 £125K",
      "製品：Nespresso対応抹茶ポッド ＋ RTDスパークリング缶",
      "売上：£8.2M（2024年実績）→ £30M（2025年予測）",
      "バリュエーション：£140M（2025年）",
      "小売：Tesco 1,200店舗・Sainsbury's・Waitrose",
      "B2B：売上の60%（Caffe Nero・Joe & The Juice等）",
      "流通：50カ国 30,000+拠点",
      "投資家：Steven Bartlett（Dragons' Den）",
      "",
      "最重要インサイト：彼らは英国のみ。",
      "北米には対抗できる競合が存在しない。",
      "私たちは日本産調達の優位性を持つ北米版PerfectTedを構築する。",
    ],
  },
  # 17. 競合ポジショニング
  {
    "type": "content",
    "tag": "COMPETITIVE LANDSCAPE",
    "title": "競合ポジショニングマップ",
    "lead": "品質（縦軸）× 利便性（横軸）",
    "body": [
      "【左上】高品質 × 低利便性",
      "  → 一保堂（品質最高、ポッドなし、在庫切れ常態化）",
      "",
      "【右下】低品質 × 中利便性",
      "  → Jade Leaf / Chamberlain Coffee（Amazon流通あり、品質に課題）",
      "",
      "【左下】低品質 × 低利便性",
      "  → Amazon汎用品（低品質、消費者の信頼なし）",
      "",
      "【右上】高品質 × 高利便性 ← 私たちのポジション（ホワイトスペース）",
      "  → PerfectTedは右上に存在するが英国のみ",
      "  → [COMPANY]は北米でこのポジションを独占する",
    ],
  },
  # 18. 競合比較表
  {
    "type": "table",
    "tag": "COMPETITIVE ANALYSIS",
    "title": "競合詳細比較",
    "thead": ["ブランド", "日本産", "ポッド", "サブスク", "価格/杯", "北米展開", "農園透明性"],
    "rows": [
      ["一保堂",            "✓", "✗", "✗", "$1.50+", "限定的", "✗"],
      ["Jade Leaf",         "✓", "✗", "✓", "$0.67", "✓", "✗"],
      ["PerfectTed",        "△", "✓", "✓", "$2.00", "英国のみ", "✗"],
      ["Chamberlain Coffee","△", "✗", "✓", "$0.85", "✓", "✗"],
      ["Cuzen Matcha",      "✓", "✗", "✓", "$1.00", "✓", "△"],
      ["[COMPANY]",         "✓✓", "✓", "✓", "$2.17", "✓", "✓✓"],
    ],
    "note": "[COMPANY]のみが全項目でチェック済み。北米唯一の高品質ポッドブランド"
  },
  # 19. 比較企業・バリュエーション
  {
    "type": "table",
    "tag": "COMPARABLE COMPANIES",
    "title": "比較企業とバリュエーション実績",
    "thead": ["企業", "調達額", "売上", "バリュエーション", "備考"],
    "rows": [
      ["PerfectTed（英）",   "£9M+",   "£8.2M", "£140M (17x rev)", "最重要コンプ。抹茶ポッドモデル"],
      ["Blank Street Coffee","$135M",  "$149M", "$500M (3.4x rev)", "抹茶が売上50%のカフェチェーン"],
      ["Jade Leaf",          "PE買収", "~$10M", "3〜5x rev推定",    "Amazon #1 → PE Exit"],
      ["RYZE Mushroom Coffee","$9M",   "$300M+","未公開",           "類似サブスク飲料。最重要参考"],
      ["AG1 (Athletic Green)","$115M", "$600M+","未公開",           "プレミアムウェルネスサブスク"],
      ["Cuzen Matcha",       "$6.8M",  "~$1M",  "未公開",           "ハードウェアは遅い→ポッドの優位性"],
    ],
    "note": "Year 3（$14M ARR）で Series A バリュエーション目標：$140M（10x rev）"
  },
  # 20. 収益モデル
  {
    "type": "threeCol",
    "tag": "BUSINESS MODEL",
    "title": "3つの収益ストリーム",
    "cols": [
      ("PRIMARY\nサブスクリプション\n売上の70%",
       "月次ポッドBOX\nスターター $35/月（15ポッド）\nレギュラー $65/月（30ポッド）\nプレミアム $95/月（45ポッド+粉末）"),
      ("SECONDARY\nB2B フードサービス\n売上の20%",
       "Nespresso対応機器を持つカフェへ\nバルクポッド供給\nASP $200〜300/月/アカウント\nカフェ訓練コスト：ゼロ"),
      ("TERTIARY\n小売・Amazon\n売上の10%",
       "Amazon（発見チャネル）\nYear2〜：Whole Foods / Target\n段階的な小売展開\n高マージンDTCを基盤にした交渉力"),
    ],
  },
  # 21. サブスク経済
  {
    "type": "bigstat",
    "tag": "UNIT ECONOMICS",
    "title": "サブスクリプション ユニットエコノミクス",
    "stats": [
      ("22x", "LTV / CAC（業界ベンチマーク：8〜15x）"),
      ("0.6ヶ月", "投資回収期間"),
      ("63%", "グロスマージン（ポッドボックス）"),
    ],
    "body": [
      "CAC（DTC混合、ソーシャル＋インフルエンサー＋リファラル）：$35",
      "月次ARPU（レギュラープラン平均）：$65",
      "LTV（12ヶ月・チャーン調整済み）：$780",
      "月次チャーン率仮定：5%（プレミアム飲料サブスクのベンチマーク）",
    ],
  },
  # 22. B2B機会
  {
    "type": "content",
    "tag": "B2B OPPORTUNITY",
    "title": "B2B フードサービス機会",
    "lead": "北米の抹茶消費の80%はカフェで行われている",
    "body": [
      "米国カフェ数：200,000+。うち60,000+がNespresso互換機を所持",
      "カフェの課題：抹茶追加には専用訓練・新機材・安定した調達が必要",
      "私たちの解決：マシンは既にある → ポッドを買うだけ → 訓練コストゼロ",
      "",
      "B2B 経済性：",
      "・ASP：$200/月/カフェアカウント（100ポッド × $2）",
      "・グロスマージン：45〜50%（マーケティングコストなし）",
      "・Year 2 目標：1,000アカウント → $2.4M 追加ARR",
      "・B2B チャーン：DTCより大幅に低い（月次発注の習慣定着）",
    ],
  },
  # 23. GTM戦略
  {
    "type": "table",
    "tag": "GO-TO-MARKET",
    "title": "フェーズ別 GTM 戦略",
    "thead": ["フェーズ", "期間", "内容", "目標"],
    "rows": [
      ["Phase 1\nDTCローンチ",   "月1〜6",  "Shopify + Amazon FBA\nマイクロインフルエンサー50名\nTikTok・Reddit コミュニティ構築", "500 サブスクライバー"],
      ["Phase 2\nB2B展開",       "月7〜18", "NYC・LA・SF・シカゴ・シアトルで\nカフェパイロット開始\nWhole Foods 商談開始",             "$3.9M ARR\n500 B2B アカウント"],
      ["Phase 3\n全国＋海外",    "Year 2〜3","Whole Foods / Target 上陸\n英国・カナダ・オーストラリア展開\n日本ローンチ検討",              "$14M ARR\n国際展開"],
    ],
  },
  # 24. マーケティング戦略
  {
    "type": "content",
    "tag": "MARKETING",
    "title": "マーケティング戦略",
    "lead": "インフルエンサーファースト → コミュニティ → 信頼 → メガ規模",
    "body": [
      "【インフルエンサー戦略：メガ < マイクロ】",
      "・50名のマイクロインフルエンサー（フォロワー5万〜50万）。エンゲージ率高・成本1/10",
      "・健康・ウェルネス・生産性ニッチ。コーヒー→抹茶切り替えナラティブが最もバイラル",
      "",
      "【コンテンツ戦略：教育が最強の武器】",
      "・「なぜあなたの抹茶は苦いのか」（信頼構築・競合破壊）",
      "・「農家からポッドへ」シリーズ（日本農家のストーリー）",
      "・科学コンテンツ：L-テアニン研究を分かりやすく解説",
      "・比較：「$2.17 vs Starbucks の $7 抹茶ラテ」",
      "",
      "【コミュニティ戦略：Redditとディスコード】",
      "・r/matchapods コミュニティを創設",
      "・サブスクライバー専用 Discord（限定コンテンツ・早期新製品・農園アップデート）",
      "・目標：Year 2 で新規獲得の30%がリファラル",
    ],
  },
  # 25. CAC計算
  {
    "type": "bigstat",
    "tag": "CUSTOMER ACQUISITION",
    "title": "Year 1 顧客獲得計算",
    "stats": [
      ("$400K", "Year 1 マーケティング予算"),
      ("$35", "ブレンドCAC 目標"),
      ("2,000", "Year 1 末 アクティブサブスクライバー目標"),
    ],
    "body": [
      "マーケティング$400K → 新規5,000名獲得（5%月次チャーンで年末2,000名維持）",
      "オーガニック・リファラルから追加 1,000名",
      "Year 1 ARR：2,000 × $65 × 12ヶ月 = $1.56M（DTCサブスクのみ）",
    ],
  },
  # 26. 日本サプライチェーン
  {
    "type": "table",
    "tag": "SUPPLY CHAIN",
    "title": "日本サプライチェーン詳細",
    "thead": ["産地", "特徴", "調達グレード", "戦略的重要度"],
    "rows": [
      ["宇治（京都）",  "800年の権威。霧と寒暖差が最高品質を生む", "セレモニアル特A", "★★★★★"],
      ["西尾（愛知）",  "全国抹茶生産量1位（20〜30%）。抹茶特化型産地",  "セレモニアル〜プレミアム",  "★★★★☆"],
      ["鹿児島",       "2024年荒茶生産量全国1位。機械化・大規模化が進む",  "デイリー〜プレミアム",  "★★★★☆"],
      ["八女（福岡）",  "全国品評会18年連続産地賞。山間の濃厚な旨味",    "超プレミアム",  "★★★☆☆"],
    ],
    "note": "産地多様化でリスク分散。宇治産をフラッグシップ、鹿児島産をボリューム供給に設計"
  },
  # 27. サプライチェーンが堀
  {
    "type": "content",
    "tag": "SUPPLY CHAIN MOAT",
    "title": "なぜサプライチェーンが最も強固な堀か",
    "lead": "抹茶において、ブランドより先にサプライチェーンを解決した者が勝つ",
    "body": [
      "競合がお金でサプライチェーンを「買う」ことはできない：",
      "   ・農家との信頼関係構築に3〜5年かかる",
      "   ・プレミアム割当は実績と安定発注量が前提",
      "   ・全契約書が日本語 → 言語・文化的参入障壁",
      "",
      "私たちの戦略：",
      "   ・今すぐ、危機の最中に農家関係を構築",
      "   ・複数年購入契約でロック価格を確保（265%高騰から絶縁）",
      "   ・農家の次世代支援プログラム（Farmer Forward）",
      "",
      "Starbucksモデルとの類比：",
      "   Starbucksが垂直統合的なコーヒー豆調達で競合を30年先行したように、",
      "   私たちは日本産抹茶で同じ堀を今から構築する",
    ],
  },
  # 28. ESG
  {
    "type": "content",
    "tag": "ESG",
    "title": "Farmer Forward：日本の茶農家を支える",
    "lead": "ESGは私たちの最強のマーケティング資産であり、供給安定の保険でもある",
    "body": [
      "【問題】",
      "・日本の茶農家人口：2010年以降40%減少",
      "・平均年齢：67歳以上。後継者危機",
      "",
      "【Farmer Forward プログラム】",
      "・市場価格以上での買取（農家に財務安定性を提供）",
      "・複数年契約（農家が投資計画を立てられる）",
      "・農家のプロフィール・ストーリーをパッケージ・ウェブサイトに掲載",
      "・地元農業大学との連携で後継者育成を支援",
      "・売上の1%を農村コミュニティ支援に充当",
      "",
      "【投資家へのリターン】",
      "・PR・メディア価値（サステナブルブランドとして差別化）",
      "・チャーン低下（ミッション志向のサブスクライバーは長期継続）",
      "・EU規制適合（欧州展開時のアドバンテージ）",
    ],
  },
  # 29. 5年財務予測
  {
    "type": "table",
    "tag": "FINANCIALS",
    "title": "5年間財務予測",
    "thead": ["", "Year 1", "Year 2", "Year 3", "Year 4", "Year 5"],
    "rows": [
      ["売上高",      "$1.2M",  "$4.8M",  "$14.0M", "$32.0M", "$65.0M"],
      ["グロスマージン","52%",   "58%",    "63%",    "66%",    "68%"],
      ["グロス利益",  "$0.6M",  "$2.8M",  "$8.8M",  "$21.1M", "$44.2M"],
      ["EBITDA",     "-$1.8M", "-$0.6M", "+$3.1M", "+$9.4M", "+$18M"],
      ["アクティブサブスク","2,000","8,000","25,000","55,000","100,000+"],
      ["B2B アカウント","200",  "1,000",  "2,500",  "5,000",  "10,000+"],
    ],
    "note": "Year 3（$14M）でEBITDAプラス転換。サブスク比率 65%+ で SaaS 並みの評価倍率"
  },
  # 30. 収益構成
  {
    "type": "bigstat",
    "tag": "REVENUE MIX",
    "title": "Year 3 収益構成（$14M）",
    "stats": [
      ("65%", "サブスクリプション（$9.1M）高 LTV・予測可能"),
      ("22%", "B2B フードサービス（$3.1M）低マーケコスト"),
      ("13%", "小売・Amazon（$1.8M）ブランド認知・発見チャネル"),
    ],
    "body": [
      "65%以上の経常収益 → SaaS並みのバリュエーション倍率を正当化",
      "B2B収益はチャーンが極低（月次発注習慣化で継続率95%+）",
    ],
  },
  # 31. 収益化への道筋
  {
    "type": "content",
    "tag": "PATH TO PROFITABILITY",
    "title": "収益化への道筋",
    "lead": "ユニット経済性は初日から黒字（LTV/CAC = 22x）",
    "body": [
      "・会社全体の EBITDA 黒字転換：月18〜24（$3M Seed 使用後）",
      "・Year 3 $14M 売上：EBITDA +$3.1M（マージン 22%）",
      "",
      "【利益率改善の3レバー】",
      "① OEM 製造コスト逓減：発注量増加でポッドCOGS が $0.85 → $0.55 へ",
      "② マーケティング効率向上：オーガニック・リファラルが拡大し CAC が $35 → $20 へ",
      "③ B2B の高貢献：マーケコスト0のB2B売上がミックスに占める比率上昇",
      "",
      "【不採算顧客を抱えない設計】",
      "サブスクモデルはユニット経済性を事前に把握可能。",
      "LTV/CAC 22x = 投資した$1が$22のリターンを生む",
    ],
  },
  # 32. 資金使途
  {
    "type": "table",
    "tag": "USE OF FUNDS",
    "title": "資金使途：$3M シードラウンド",
    "thead": ["カテゴリ", "金額", "比率", "内訳"],
    "rows": [
      ["製品開発・OEM 設立",  "$800K", "27%", "カプセル処方・テスト $200K\nOEM工場初回セットアップ $400K\nISO認証・品質テスト $200K"],
      ["在庫・サプライチェーン","$700K","23%", "農家との前払い契約 $400K\n倉庫設立・6ヶ月在庫バッファ $300K"],
      ["マーケティング・成長",  "$800K", "27%", "インフルエンサー・コンテンツ $300K\n有料ソーシャル（TikTok・IG）$300K\nPR・コミュニティ構築 $200K"],
      ["チーム採用",           "$500K", "17%", "サプライチェーン責任者 $180K\n成長・マーケ責任者 $150K\nCS リード $100K・その他 $70K"],
      ["オペレーション・法務",  "$200K",  "7%", "法務（IP・契約・サブスク利用規約）$80K\nシステム・コンプライアンス $120K"],
    ],
  },
  # 33. マイルストーン
  {
    "type": "table",
    "tag": "MILESTONES",
    "title": "マイルストーン：$3M で達成すること",
    "thead": ["時期", "マイルストーン", "目標指標"],
    "rows": [
      ["Month 1〜3",  "OEM確定・農家契約締結\nShopify＋サブスクプラットフォーム構築\nインフルエンサー50名シーディング",  "農家MOU締結"],
      ["Month 3",    "製品ローンチ（Shopify DTC + Amazon FBA）",  "初販売"],
      ["Month 6",    "B2B パイロット開始（5都市）",  "500 サブスクライバー\n$390K ARR"],
      ["Month 12",   "Amazon 抹茶ポッドカテゴリ TOP 5\nWhole Foods 商談開始",  "2,000 サブスク\n$1.56M ARR"],
      ["Month 18",   "Series A 準備完了",  "5,000 サブスク\n$3.9M ARR\n500 B2B アカウント"],
    ],
    "note": "Series A 目標：$8〜12M を $40〜60M pre-money で調達"
  },
  # 34. KPI
  {
    "type": "table",
    "tag": "KEY METRICS",
    "title": "KPI ダッシュボード",
    "thead": ["KPI", "Year 1 目標", "Year 3 目標", "ベンチマーク"],
    "rows": [
      ["アクティブサブスクライバー数", "2,000",  "25,000",  "—"],
      ["MRR / ARR",                   "$130K / $1.56M", "$950K / $11.4M", "—"],
      ["月次グロスチャーン率",         "< 5%",   "< 4%",   "業界平均 7〜10%"],
      ["Net Revenue Retention",       "> 105%", "> 115%", "優秀 SaaS > 110%"],
      ["ブレンドCAC",                  "< $40",  "< $20",  "—"],
      ["LTV / CAC",                   "> 15x",  "> 30x",  "業界ベンチマーク 8〜15x"],
      ["B2B アカウント数",             "200",    "2,500",  "—"],
      ["グロスマージン",               "52%",    "63%",    "DTC飲料 平均 50〜60%"],
    ],
  },
  # 35. ビジョン
  {
    "type": "content",
    "tag": "VISION",
    "title": "ビジョン：抹茶を超えたプラットフォームへ",
    "lead": "Year 5+：「日本産ウェルネス飲料の Nespresso」",
    "body": [
      "【Year 1〜3：北米 No.1 抹茶サブスクブランドの確立】",
      "・Nespresso対応抹茶ポッドで北米カテゴリを独占",
      "",
      "【Year 3〜5：製品ライン拡張】",
      "・ほうじ茶ポッド・玉露ポッド・アダプトゲン×抹茶ブレンド",
      "・機能性ラインナップ（抹茶 + Lion's Mane + L-テアニン）",
      "",
      "【Year 5+：プラットフォーム化】",
      "・日本産プレミアム原料の調達プラットフォームとしてライセンス提供",
      "・カフェ向けポッドのホワイトラベル",
      "・日本産プレミアム食材の西洋向け総合ゲートウェイ",
      "",
      "【長期Exit シナリオ】",
      "・Nestle（Nespresso親会社）による戦略的買収",
      "・Keurig Dr Pepper・AB InBev・伊藤園・サントリー",
      "・$200M+ 売上到達後の IPO パス（AG1 モデル）",
    ],
  },
  # 36. 5つの不公平な優位性
  {
    "type": "content",
    "tag": "UNFAIR ADVANTAGES",
    "title": "5つの不公平な優位性",
    "lead": "これらは短期間では複製できない構造的優位性",
    "body": [
      "① タイミング：需要は爆発、供給は追いつかない。今が先行者優位を獲得する最後のウィンドウ",
      "",
      "② フォーマット革新：ポッドは採用の最大障壁（準備の複雑さ）を完全に排除する",
      "",
      "③ サプライチェーン堀：日本農家との関係は数年かけて構築される。お金では買えない",
      "",
      "④ 既存インフラ活用：3,000万台のNespressoマシン = ハードウェアリスクゼロ",
      "",
      "⑤ 実証済みモデル：PerfectTedが同一モデルを英国で成功実証済み。私たちはより大きな市場で再現する",
    ],
  },
  # 37. 調達内容
  {
    "type": "ask",
    "tag": "THE ASK",
    "title": "調達内容",
    "stats": [
      ("$3M", "シードラウンド調達額"),
      ("$12M", "プレマネーバリュエーション"),
      ("18ヶ月", "ランウェイ → $4.8M ARR → Series A"),
    ],
    "body": [
      "リード投資家：$1.5M 募集中",
      "クロージング目標：2026年Q2",
      "投資形態：SAFE または プライスドラウンド（優先株）",
      "ボード構成：創業者コントロール ＋ 投資家1席",
    ],
  },
  # 38. チーム
  {
    "type": "content",
    "tag": "TEAM",
    "title": "チーム",
    "lead": "日本の文化的真正性 × シリコンバレーの流通ノウハウ",
    "body": [
      "CEO / Co-Founder：[氏名] — [バックグラウンド]",
      "COO / Co-Founder：[氏名] — サプライチェーン / 日本オペレーション経験",
      "",
      "【アドバイザー】",
      "・日本茶業界エキスパート：[氏名]",
      "・CPG / 飲料業界エグゼクティブ：[氏名]",
      "・DTCサブスクリプション専門家：[氏名]",
      "",
      "【求めるもの（チームビルディング中）】",
      "・Head of Supply Chain / Japan Operations",
      "・Head of Growth / Marketing（DTC サブスク経験必須）",
      "・Customer Experience Lead",
      "",
      "※ このスライドは最も重要。投資判断の決め手になる。",
      "  投資家ミーティング前に実名・具体的経歴を必ず記載すること。",
    ],
  },
  # 39. まとめ
  {
    "type": "closing",
    "tag": "CLOSING",
    "title": "今こそ、私たちこそ",
    "body": [
      "✓  市場：$4.3B → $8.9B（CAGR 7.8%）。北米は CAGR 11% で最速成長",
      "✓  競合：北米の抹茶ポッド市場は競合ゼロ",
      "✓  証明：PerfectTed が同一モデルで CAGR 532%・£140M バリュエーションを達成",
      "✓  緊急性：供給危機により先行者優位は今しか獲得できない",
      "✓  経済性：サブスクリプションで LTV/CAC 22x・予測可能な経常収益",
      "",
      "私たちは抹茶をコーヒーと並ぶカフェイン選択肢として確立する。",
      "その入口は、既に3,000万人の家にある机の上にある。",
      "",
      "《 The Nespresso of Matcha 》",
    ],
  },
]

# ─────────────────────────────────────────────────────────────
# PPTX ヘルパー関数
# ─────────────────────────────────────────────────────────────

def add_textbox(slide, text, left, top, width, height,
                font_size=18, bold=False, italic=False,
                color=None, align=PP_ALIGN.LEFT,
                font_name="Hiragino Sans", wrap=True):
    if color is None:
        color = DTXT
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_name
    return txBox

def add_rect(slide, left, top, width, height, fill_color,
             line_color=None, line_width=0):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape

def set_dark_gradient_bg(slide):
    """OOXML injection でダークグラデーション背景を設定"""
    cSld = slide._element.cSld
    old = cSld.find(qn('p:bg'))
    if old is not None:
        cSld.remove(old)
    bg_xml = (
        '<p:bg xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"'
        ' xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        '<p:bgPr>'
        '<a:gradFill rotWithShape="1">'
        '<a:gsLst>'
        '<a:gs pos="0"><a:srgbClr val="1A2B42"/></a:gs>'
        '<a:gs pos="100000"><a:srgbClr val="050911"/></a:gs>'
        '</a:gsLst>'
        '<a:lin ang="8100000" scaled="0"/>'
        '</a:gradFill>'
        '<a:effectLst/>'
        '</p:bgPr>'
        '</p:bg>'
    )
    cSld.insert(0, etree.fromstring(bg_xml))

def set_light_bg(slide):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = LBKG

def add_dark_decos(slide):
    """ダークスライドに装飾的な大円を追加"""
    # 大きな円（右上）
    r1 = Cm(9.5)
    sp1 = slide.shapes.add_shape(9, SW - r1 * 1.4, -r1 * 0.5, r1 * 2, r1 * 2)
    sp1.fill.solid()
    sp1.fill.fore_color.rgb = DECO_G
    sp1.line.fill.background()
    # 小さな円（左下）
    r2 = Cm(5.5)
    sp2 = slide.shapes.add_shape(9, -r2 * 0.5, SH - r2 * 1.2, r2 * 2, r2 * 2)
    sp2.fill.solid()
    sp2.fill.fore_color.rgb = DECO_B
    sp2.line.fill.background()

def add_light_header(slide, tag, title):
    """ライトスライド用ミニマルヘッダー"""
    # 最上部アクセントストリップ（マッチャグリーン）
    add_rect(slide, Cm(0), Cm(0), SW, Cm(0.13), GREEN)
    # 左縦アクセントライン
    add_rect(slide, Cm(0), Cm(0.13), Cm(0.22), Cm(2.4), GREEN)
    # タグテキスト
    if tag:
        add_textbox(slide, tag,
                    Cm(0.6), Cm(0.22), Cm(20), Cm(0.55),
                    font_size=8, bold=True, color=GDARK,
                    align=PP_ALIGN.LEFT)
    # タイトルテキスト
    add_textbox(slide, title,
                Cm(0.6), Cm(0.8), Cm(32.5), Cm(1.6),
                font_size=23, bold=True, color=DTXT,
                align=PP_ALIGN.LEFT)
    # セパレーターライン
    add_rect(slide, Cm(0.6), Cm(2.5), Cm(32.5), Cm(0.05), BRDR)

def add_footer_light(slide, slide_num, total):
    add_rect(slide, Cm(0.6), Cm(18.4), Cm(32.5), Cm(0.05), BRDR)
    add_textbox(slide, f"{slide_num} / {total}",
                Cm(30), Cm(18.5), Cm(3.5), Cm(0.5),
                font_size=8, color=MTXT, align=PP_ALIGN.RIGHT)
    add_textbox(slide, "[COMPANY]  ·  Confidential",
                Cm(0.6), Cm(18.5), Cm(15), Cm(0.5),
                font_size=8, color=MTXT, align=PP_ALIGN.LEFT)

def add_footer_dark(slide, slide_num, total):
    add_rect(slide, Cm(0.6), Cm(18.4), Cm(32.5), Cm(0.05), DECO_B)
    add_textbox(slide, f"{slide_num} / {total}",
                Cm(30), Cm(18.5), Cm(3.5), Cm(0.5),
                font_size=8, color=DIM, align=PP_ALIGN.RIGHT)
    add_textbox(slide, "[COMPANY]  ·  Confidential  ·  2026",
                Cm(0.6), Cm(18.5), Cm(18), Cm(0.5),
                font_size=8, color=DIM, align=PP_ALIGN.LEFT)

# ─────────────────────────────────────────────────────────────
# スライドタイプ別レンダリング関数
# ─────────────────────────────────────────────────────────────

def render_cover(slide, prs, data):
    set_dark_gradient_bg(slide)
    add_dark_decos(slide)
    # 左縦アクセントバー（マッチャグリーン）
    add_rect(slide, Cm(0), Cm(0), Cm(0.55), SH, GREEN)
    # 会社名（超大）
    add_textbox(slide, data["title"],
                Cm(1.5), Cm(1.0), Cm(28), Cm(2.8),
                font_size=64, bold=True, color=WHITE,
                align=PP_ALIGN.LEFT)
    # タグライン
    add_textbox(slide, data["subtitle"],
                Cm(1.5), Cm(3.9), Cm(26), Cm(1.2),
                font_size=24, bold=False, color=GREEN,
                align=PP_ALIGN.LEFT)
    # サブタイトル下のアクセントライン
    add_rect(slide, Cm(1.5), Cm(5.2), Cm(22), Cm(0.06), GREEN)
    # ボディテキスト
    body_text = "\n".join(data.get("body", []))
    add_textbox(slide, body_text,
                Cm(1.5), Cm(5.5), Cm(27), Cm(9.5),
                font_size=13, color=DIM,
                align=PP_ALIGN.LEFT)
    # Seed Roundバッジ（右下）
    add_rect(slide, Cm(23.5), Cm(16.5), Cm(9.8), Cm(0.85),
             RGBColor(0x14, 0x3D, 0x27))
    add_textbox(slide, "SEED ROUND  ·  $3,000,000  ·  2026",
                Cm(23.7), Cm(16.58), Cm(9.4), Cm(0.72),
                font_size=9, color=GREEN, align=PP_ALIGN.CENTER)
    add_footer_dark(slide, 1, len(SLIDES))

def render_content(slide, prs, data, slide_num):
    set_light_bg(slide)
    add_light_header(slide, data.get("tag", ""), data.get("title", ""))
    add_footer_light(slide, slide_num, len(SLIDES))
    y = Cm(2.75)
    if data.get("lead"):
        add_rect(slide, Cm(0.6), y, Cm(0.22), Cm(0.95), GREEN)
        add_textbox(slide, data["lead"],
                    Cm(1.05), y + Cm(0.05), Cm(31.5), Cm(1.0),
                    font_size=14, bold=True, color=GDARK,
                    align=PP_ALIGN.LEFT)
        y += Cm(1.2)
    body_text = "\n".join(data.get("body", []))
    add_textbox(slide, body_text,
                Cm(0.8), y, Cm(32), SH - y - Cm(1.0),
                font_size=12, color=DTXT,
                align=PP_ALIGN.LEFT)

def render_bigstat(slide, prs, data, slide_num):
    set_light_bg(slide)
    add_light_header(slide, data.get("tag", ""), data.get("title", ""))
    add_footer_light(slide, slide_num, len(SLIDES))
    stats = data.get("stats", [])
    n = len(stats)
    gap = Cm(0.45)
    card_w = (Cm(32.5) - gap * (n - 1)) / n
    card_colors = [GREEN, BLUE_H, GOLD]
    card_bgs    = [CARD_G, CARD_B, CARD_A]
    card_top = Cm(2.75)
    card_h   = Cm(5.6)
    for i, (val, lbl) in enumerate(stats):
        lx = Cm(0.55) + i * (card_w + gap)
        add_rect(slide, lx, card_top, card_w, card_h, card_bgs[i % 3])
        # 上部カラーボーダー
        add_rect(slide, lx, card_top, card_w, Cm(0.2), card_colors[i % 3])
        # 大きな統計数値
        add_textbox(slide, val,
                    lx, card_top + Cm(0.3), card_w, Cm(2.6),
                    font_size=52, bold=True, color=card_colors[i % 3],
                    align=PP_ALIGN.CENTER)
        # ラベル
        add_textbox(slide, lbl,
                    lx + Cm(0.2), card_top + Cm(3.1), card_w - Cm(0.4), Cm(2.3),
                    font_size=11, color=DTXT,
                    align=PP_ALIGN.CENTER)
    body_text = "\n".join(data.get("body", []))
    if body_text:
        add_textbox(slide, body_text,
                    Cm(0.8), card_top + card_h + Cm(0.4),
                    Cm(32), SH - (card_top + card_h + Cm(0.4)) - Cm(1.0),
                    font_size=12, color=DTXT,
                    align=PP_ALIGN.LEFT)

def render_table(slide, prs, data, slide_num):
    set_light_bg(slide)
    add_light_header(slide, data.get("tag", ""), data.get("title", ""))
    add_footer_light(slide, slide_num, len(SLIDES))
    thead = data.get("thead", [])
    rows  = data.get("rows", [])
    note  = data.get("note", "")
    ncols = len(thead)
    nrows = len(rows)
    if ncols == 0 or nrows == 0:
        return
    table_left = Cm(0.55)
    table_top  = Cm(2.75)
    table_w    = Cm(32.5)
    note_h     = Cm(1.0) if note else Cm(0.2)
    table_h    = SH - table_top - Cm(1.2) - note_h
    row_h      = table_h / (nrows + 1)
    col_w      = table_w / ncols
    # ヘッダー行（ダークネイビー）
    add_rect(slide, table_left, table_top, table_w, row_h, DTXT)
    for ci, th in enumerate(thead):
        add_textbox(slide, th,
                    table_left + col_w * ci + Cm(0.2), table_top + Cm(0.1),
                    col_w - Cm(0.3), row_h - Cm(0.15),
                    font_size=11, bold=True, color=WHITE,
                    align=PP_ALIGN.LEFT)
    # データ行
    for ri, row in enumerate(rows):
        bg = CARD_G if ri % 2 == 0 else LBKG
        ry = table_top + row_h * (ri + 1)
        add_rect(slide, table_left, ry, table_w, row_h, bg)
        if ri % 2 == 0:
            add_rect(slide, table_left, ry, Cm(0.13), row_h, GREEN)
        for ci, cell in enumerate(row):
            add_textbox(slide, str(cell),
                        table_left + col_w * ci + Cm(0.2), ry + Cm(0.08),
                        col_w - Cm(0.3), row_h - Cm(0.12),
                        font_size=11, color=DTXT,
                        align=PP_ALIGN.LEFT)
    if note:
        ny = table_top + row_h * (nrows + 1) + Cm(0.15)
        add_textbox(slide, f"※ {note}",
                    table_left, ny, table_w, note_h,
                    font_size=9, color=MTXT,
                    align=PP_ALIGN.LEFT)

def render_three_col(slide, prs, data, slide_num):
    set_light_bg(slide)
    add_light_header(slide, data.get("tag", ""), data.get("title", ""))
    add_footer_light(slide, slide_num, len(SLIDES))
    cols = data.get("cols", [])
    n = len(cols)
    gap = Cm(0.45)
    card_w = (Cm(32.5) - gap * (n - 1)) / n
    hdr_colors = [GREEN, BLUE_H, GOLD]
    bg_colors  = [CARD_G, CARD_B, CARD_A]
    card_top = Cm(2.75)
    card_h   = SH - card_top - Cm(1.2)
    hdr_h    = Cm(2.3)
    for i, (hdr, body) in enumerate(cols):
        lx = Cm(0.55) + i * (card_w + gap)
        add_rect(slide, lx, card_top, card_w, card_h, bg_colors[i % 3])
        add_rect(slide, lx, card_top, card_w, hdr_h, hdr_colors[i % 3])
        add_textbox(slide, hdr,
                    lx + Cm(0.2), card_top + Cm(0.2), card_w - Cm(0.4), hdr_h - Cm(0.3),
                    font_size=13, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER)
        add_textbox(slide, body,
                    lx + Cm(0.3), card_top + hdr_h + Cm(0.25), card_w - Cm(0.5), card_h - hdr_h - Cm(0.5),
                    font_size=11, color=DTXT,
                    align=PP_ALIGN.LEFT)

def render_three_step(slide, prs, data, slide_num):
    set_light_bg(slide)
    add_light_header(slide, data.get("tag", ""), data.get("title", ""))
    add_footer_light(slide, slide_num, len(SLIDES))
    steps = data.get("steps", [])
    n = len(steps)
    step_colors = [GREEN, BLUE_H, GOLD]
    step_w = Cm(32.5) / n
    # コネクターライン
    add_rect(slide, Cm(4), Cm(7.5), Cm(25.5), Cm(0.06), BRDR)
    for i, (hdr, body) in enumerate(steps):
        lx = Cm(0.55) + i * step_w
        cx = lx + step_w / 2
        r  = Cm(1.5)
        # 円
        sp = slide.shapes.add_shape(9, cx - r, Cm(5.2), r * 2, r * 2)
        sp.fill.solid()
        sp.fill.fore_color.rgb = step_colors[i % 3]
        sp.line.fill.background()
        # 番号
        add_textbox(slide, str(i + 1),
                    cx - r, Cm(5.5), r * 2, r * 1.8,
                    font_size=20, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER)
        # ヘッダー
        add_textbox(slide, hdr,
                    lx + Cm(0.3), Cm(8.3), step_w - Cm(0.6), Cm(2.2),
                    font_size=13, bold=True, color=step_colors[i % 3],
                    align=PP_ALIGN.CENTER)
        # ボディ
        add_textbox(slide, body,
                    lx + Cm(0.3), Cm(10.5), step_w - Cm(0.6), Cm(7.5),
                    font_size=12, color=DTXT,
                    align=PP_ALIGN.CENTER)

def render_ask(slide, prs, data, slide_num):
    set_dark_gradient_bg(slide)
    add_dark_decos(slide)
    # 左アクセント（ゴールド）
    add_rect(slide, Cm(0), Cm(0), Cm(0.55), SH, GOLD)
    add_textbox(slide, data.get("tag", ""),
                Cm(1.3), Cm(0.35), Cm(12), Cm(0.7),
                font_size=9, bold=True, color=GOLD,
                align=PP_ALIGN.LEFT)
    add_textbox(slide, data.get("title", ""),
                Cm(1.3), Cm(1.0), Cm(30), Cm(1.8),
                font_size=34, bold=True, color=WHITE,
                align=PP_ALIGN.LEFT)
    add_rect(slide, Cm(1.3), Cm(3.0), Cm(31.5), Cm(0.07), GOLD)
    stats = data.get("stats", [])
    n = len(stats)
    card_w = Cm(30.5 / n)
    for i, (val, lbl) in enumerate(stats):
        lx = Cm(1.3) + i * (card_w + Cm(0.45))
        add_rect(slide, lx, Cm(3.2), card_w, Cm(4.8), DECO_G)
        add_rect(slide, lx, Cm(3.2), card_w, Cm(0.18), GOLD)
        add_textbox(slide, val,
                    lx, Cm(3.5), card_w, Cm(2.4),
                    font_size=50, bold=True, color=GOLD,
                    align=PP_ALIGN.CENTER)
        add_textbox(slide, lbl,
                    lx + Cm(0.2), Cm(6.1), card_w - Cm(0.4), Cm(1.7),
                    font_size=11, color=GREEN,
                    align=PP_ALIGN.CENTER)
    body_text = "\n".join(data.get("body", []))
    add_textbox(slide, body_text,
                Cm(1.3), Cm(8.4), Cm(31), Cm(9.5),
                font_size=13, color=DIM,
                align=PP_ALIGN.LEFT)
    add_footer_dark(slide, slide_num, len(SLIDES))

def render_closing(slide, prs, data, slide_num):
    set_dark_gradient_bg(slide)
    add_dark_decos(slide)
    # 左アクセント（マッチャグリーン）
    add_rect(slide, Cm(0), Cm(0), Cm(0.55), SH, GREEN)
    add_textbox(slide, data.get("tag", ""),
                Cm(1.3), Cm(0.35), Cm(12), Cm(0.7),
                font_size=9, bold=True, color=GREEN,
                align=PP_ALIGN.LEFT)
    add_textbox(slide, data.get("title", ""),
                Cm(1.3), Cm(1.0), Cm(30), Cm(1.8),
                font_size=34, bold=True, color=WHITE,
                align=PP_ALIGN.LEFT)
    add_rect(slide, Cm(1.3), Cm(3.0), Cm(31.5), Cm(0.07), GREEN)
    body_text = "\n".join(data.get("body", []))
    add_textbox(slide, body_text,
                Cm(1.3), Cm(3.3), Cm(31), Cm(14.5),
                font_size=14, color=DIM,
                align=PP_ALIGN.LEFT)
    add_footer_dark(slide, slide_num, len(SLIDES))

# ─────────────────────────────────────────────────────────────
# PPTX 生成
# ─────────────────────────────────────────────────────────────

def build_pptx(output_path):
    prs = Presentation()
    prs.slide_width  = SW
    prs.slide_height = SH
    blank_layout = prs.slide_layouts[6]
    for i, data in enumerate(SLIDES):
        slide = prs.slides.add_slide(blank_layout)
        t  = data["type"]
        sn = i + 1
        if   t == "cover":     render_cover(slide, prs, data)
        elif t == "content":   render_content(slide, prs, data, sn)
        elif t == "bigstat":   render_bigstat(slide, prs, data, sn)
        elif t == "table":     render_table(slide, prs, data, sn)
        elif t == "threeCol":  render_three_col(slide, prs, data, sn)
        elif t == "threeStep": render_three_step(slide, prs, data, sn)
        elif t == "ask":       render_ask(slide, prs, data, sn)
        elif t == "closing":   render_closing(slide, prs, data, sn)
        else:                  render_content(slide, prs, data, sn)
    prs.save(output_path)
    print(f"✓ PPTX saved: {output_path}")

# ─────────────────────────────────────────────────────────────
# PDF 生成（reportlab）
# ─────────────────────────────────────────────────────────────

M = mm  # alias

def pdf_dark_bg(c, W, H):
    """ダーク背景（グラデーション風）"""
    c.setFillColor(rl_colors.HexColor('#090D18'))
    c.rect(0, 0, W, H, stroke=0, fill=1)
    # グラデーション風：上部を少し明るく
    c.setFillColor(rl_colors.HexColor('#142033'))
    c.rect(0, H * 0.55, W, H * 0.45, stroke=0, fill=1)
    # 装飾円（右上）
    c.setFillColor(rl_colors.HexColor('#0B2315'))
    c.circle(W - 40*M, H + 10*M, 90*M, stroke=0, fill=1)
    # 装飾円（左下）
    c.setFillColor(rl_colors.HexColor('#071428'))
    c.circle(0, 0, 55*M, stroke=0, fill=1)

def pdf_draw_slide(c, data, slide_num, total):
    W, H = PAGE_W, PAGE_H
    t = data["type"]

    # ── 背景 ──────────────────────────────────────────────
    if t in ("cover", "ask", "closing"):
        pdf_dark_bg(c, W, H)
    else:
        c.setFillColor(RL_LBKG)
        c.rect(0, 0, W, H, stroke=0, fill=1)

    # ── カバー ─────────────────────────────────────────────
    if t == "cover":
        # 左縦ストライプ
        c.setFillColor(RL_GREEN)
        c.rect(0, 0, 5.5*M, H, stroke=0, fill=1)
        # タイトル
        c.setFont(JP_FONT, 48)
        c.setFillColor(RL_WHITE)
        c.drawString(15*M, H - 28*M, data["title"])
        # サブタイトル
        c.setFont(JP_FONT, 22)
        c.setFillColor(RL_GREEN)
        c.drawString(15*M, H - 46*M, data["subtitle"])
        # アクセントライン
        c.setFillColor(RL_GREEN)
        c.rect(15*M, H - 53*M, 150*M, 0.6*M, stroke=0, fill=1)
        # ボディ
        c.setFont(JP_FONT, 11)
        c.setFillColor(RL_DIM)
        y = H - 62*M
        for line in data.get("body", []):
            if line:
                c.drawString(15*M, y, line)
            y -= 9.5*M
        # バッジ
        c.setFillColor(rl_colors.HexColor('#14402A'))
        c.rect(W - 115*M, 15*M, 110*M, 10*M, stroke=0, fill=1)
        c.setFont(JP_FONT, 8)
        c.setFillColor(RL_GREEN)
        badge_text = "SEED ROUND  ·  $3,000,000  ·  2026"
        bw = c.stringWidth(badge_text, JP_FONT, 8)
        c.drawString(W - 115*M + (110*M - bw)/2, 18*M, badge_text)

    # ── ASK / CLOSING ───────────────────────────────────────
    elif t in ("ask", "closing"):
        accent = RL_GOLD if t == "ask" else RL_GREEN
        # 左縦ストライプ
        c.setFillColor(accent)
        c.rect(0, 0, 5.5*M, H, stroke=0, fill=1)
        # タグ
        c.setFont(JP_FONT, 8)
        c.setFillColor(accent)
        c.drawString(13*M, H - 10*M, data.get("tag", ""))
        # タイトル
        c.setFont(JP_FONT, 28)
        c.setFillColor(RL_WHITE)
        c.drawString(13*M, H - 25*M, data.get("title", ""))
        # セパレーター
        c.setFillColor(accent)
        c.rect(13*M, H - 31*M, W - 20*M, 0.7*M, stroke=0, fill=1)
        if t == "ask":
            stats = data.get("stats", [])
            n = len(stats)
            sw = (W - 26*M) / n
            for i, (val, lbl) in enumerate(stats):
                lx = 13*M + i * sw
                c.setFillColor(rl_colors.HexColor('#0B2315'))
                c.rect(lx, H - 77*M, sw - 3*M, 44*M, stroke=0, fill=1)
                c.setFillColor(RL_GOLD)
                c.rect(lx, H - 33*M, sw - 3*M, 2*M, stroke=0, fill=1)
                c.setFont(JP_FONT, 36)
                c.setFillColor(RL_GOLD)
                vw = c.stringWidth(val, JP_FONT, 36)
                c.drawString(lx + (sw - 3*M - vw)/2, H - 60*M, val)
                c.setFont(JP_FONT, 10)
                c.setFillColor(RL_GREEN)
                lw2 = c.stringWidth(lbl, JP_FONT, 10)
                c.drawString(lx + (sw - 3*M - lw2)/2, H - 73*M, lbl)
            y0 = H - 84*M
        else:
            y0 = H - 40*M
        c.setFont(JP_FONT, 12)
        c.setFillColor(RL_DIM)
        for line in data.get("body", []):
            if line:
                c.drawString(13*M, y0, line)
            y0 -= 9.5*M

    # ── 通常スライド ────────────────────────────────────────
    else:
        # 最上部ストリップ（マッチャグリーン）
        c.setFillColor(RL_GREEN)
        c.rect(0, H - 1.3*M, W, 1.3*M, stroke=0, fill=1)
        # 左縦アクセント
        c.setFillColor(RL_GREEN)
        c.rect(0, H - 24*M, 2.2*M, 22.7*M, stroke=0, fill=1)
        # タグ
        c.setFont(JP_FONT, 7)
        c.setFillColor(RL_GDARK)
        c.drawString(5.5*M, H - 8*M, data.get("tag", ""))
        # タイトル
        c.setFont(JP_FONT, 19)
        c.setFillColor(RL_DTXT)
        c.drawString(5.5*M, H - 19*M, data.get("title", ""))
        # セパレーター
        c.setStrokeColor(RL_BRDR)
        c.setLineWidth(0.5)
        c.line(5.5*M, H - 26*M, W - 5.5*M, H - 26*M)

        y = H - 34*M

        # bigstat
        if t == "bigstat":
            stats = data.get("stats", [])
            ns = len(stats)
            bw = (W - 16*M) / ns
            colors_s = [RL_GREEN, RL_BLUE_H, RL_GOLD]
            bgs_s    = [RL_CARD_G, RL_CARD_B, RL_CARD_A]
            for i, (val, lbl) in enumerate(stats):
                lx = 8*M + i * bw
                c.setFillColor(bgs_s[i % 3])
                c.roundRect(lx, H - 90*M, bw - 4*M, 56*M, 3*M, stroke=0, fill=1)
                c.setFillColor(colors_s[i % 3])
                c.rect(lx, H - 34*M, bw - 4*M, 2*M, stroke=0, fill=1)
                c.setFont(JP_FONT, 36)
                c.setFillColor(colors_s[i % 3])
                vw = c.stringWidth(val, JP_FONT, 36)
                c.drawString(lx + (bw - 4*M - vw)/2, H - 62*M, val)
                c.setFont(JP_FONT, 9)
                c.setFillColor(RL_DTXT)
                for li, ln in enumerate(lbl.split('\n')):
                    lw2 = c.stringWidth(ln, JP_FONT, 9)
                    c.drawString(lx + (bw - 4*M - lw2)/2, H - 76*M - li*9*M, ln)
            y = H - 96*M
            c.setFont(JP_FONT, 11)
            c.setFillColor(RL_DTXT)
            for line in data.get("body", []):
                if line:
                    c.drawString(8*M, y, line)
                y -= 9*M

        # table
        elif t == "table":
            thead = data.get("thead", [])
            rows  = data.get("rows", [])
            nc = len(thead)
            nr = len(rows)
            note = data.get("note", "")
            if nc > 0 and nr > 0:
                note_h = 10*M if note else 2*M
                tH = H - 34*M - 8*M - note_h
                rh = tH / (nr + 1)
                cw = (W - 11*M) / nc
                c.setFillColor(RL_DTXT)
                c.rect(5.5*M, H - 34*M - rh, W - 11*M, rh, stroke=0, fill=1)
                for ci, th in enumerate(thead):
                    c.setFont(JP_FONT, 10)
                    c.setFillColor(RL_WHITE)
                    c.drawString(5.5*M + cw*ci + 2*M, H - 34*M - rh + 3*M, str(th))
                for ri, row in enumerate(rows):
                    ry = H - 34*M - rh*(ri+2)
                    bg = RL_CARD_G if ri % 2 == 0 else RL_LBKG
                    c.setFillColor(bg)
                    c.rect(5.5*M, ry, W - 11*M, rh, stroke=0, fill=1)
                    if ri % 2 == 0:
                        c.setFillColor(RL_GREEN)
                        c.rect(5.5*M, ry, 1.3*M, rh, stroke=0, fill=1)
                    for ci, cell in enumerate(row):
                        c.setFont(JP_FONT, 9)
                        c.setFillColor(RL_DTXT)
                        lines_c = str(cell).split('\n')
                        for li, ln in enumerate(lines_c[:2]):
                            c.drawString(5.5*M + cw*ci + 2.5*M, ry + rh - 8*M - li*7*M, ln)
                if note:
                    ny = H - 34*M - rh*(nr+1) - 7*M
                    c.setFont(JP_FONT, 8)
                    c.setFillColor(RL_MTXT)
                    c.drawString(5.5*M, ny, f"※ {note}")

        # threeCol / threeStep
        elif t in ("threeCol", "threeStep"):
            items = data.get("cols", data.get("steps", []))
            ni = len(items)
            cw = (W - 21*M) / ni
            hdr_cs = [RL_GREEN, RL_BLUE_H, RL_GOLD]
            bg_cs  = [RL_CARD_G, RL_CARD_B, RL_CARD_A]
            for i, item in enumerate(items):
                hdr_t, body_t = item
                lx = 8*M + i * (cw + 2.5*M)
                c.setFillColor(bg_cs[i % 3])
                c.roundRect(lx, H - 165*M, cw, 130*M, 3*M, stroke=0, fill=1)
                c.setFillColor(hdr_cs[i % 3])
                c.roundRect(lx, H - 55*M, cw, 20*M, 3*M, stroke=0, fill=1)
                c.setFont(JP_FONT, 10)
                c.setFillColor(RL_WHITE)
                for hi, hl in enumerate(hdr_t.split('\n')):
                    hw = c.stringWidth(hl, JP_FONT, 10)
                    c.drawString(lx + (cw - hw)/2, H - 41*M - hi*10*M, hl)
                c.setFont(JP_FONT, 10)
                c.setFillColor(RL_DTXT)
                y2 = H - 60*M
                for line in body_t.split('\n'):
                    if line:
                        c.drawString(lx + 2*M, y2, line)
                    y2 -= 8.5*M

        # content（デフォルト）
        else:
            lead = data.get("lead", "")
            if lead:
                c.setFillColor(RL_GREEN)
                c.rect(5.5*M, y - 1*M, 2.2*M, 10*M, stroke=0, fill=1)
                c.setFont(JP_FONT, 12)
                c.setFillColor(RL_GDARK)
                for line in lead.split('\n'):
                    c.drawString(9.5*M, y, line)
                    y -= 10.5*M
                y -= 2*M
            c.setFont(JP_FONT, 11)
            c.setFillColor(RL_DTXT)
            for line in data.get("body", []):
                if line.startswith("【"):
                    c.setFillColor(RL_GDARK)
                elif line == "":
                    c.setFillColor(RL_DTXT)
                else:
                    c.setFillColor(RL_DTXT)
                if line:
                    c.drawString(7*M, y, line)
                y -= 9.5*M

    # ── フッター（全スライド共通）────────────────────────────
    pg_text = f"{slide_num} / {total}"
    if t in ("cover", "ask", "closing"):
        c.setStrokeColor(rl_colors.HexColor('#1A2E45'))
        c.setLineWidth(0.5)
        c.line(5.5*M, 8*M, W - 5.5*M, 8*M)
        c.setFont(JP_FONT, 7)
        c.setFillColor(RL_DIM)
        pw = c.stringWidth(pg_text, JP_FONT, 7)
        c.drawString(W - pw - 8*M, 3*M, pg_text)
        c.drawString(8*M, 3*M, "[COMPANY]  ·  Confidential  ·  2026")
    else:
        c.setStrokeColor(RL_BRDR)
        c.setLineWidth(0.5)
        c.line(5.5*M, 8*M, W - 5.5*M, 8*M)
        c.setFont(JP_FONT, 7)
        c.setFillColor(RL_MTXT)
        pw = c.stringWidth(pg_text, JP_FONT, 7)
        c.drawString(W - pw - 8*M, 3*M, pg_text)
        c.drawString(8*M, 3*M, "[COMPANY]  ·  Confidential  ·  2026")


def build_pdf(output_path):
    c = pdf_canvas.Canvas(output_path, pagesize=(PAGE_W, PAGE_H))
    total = len(SLIDES)
    for i, data in enumerate(SLIDES):
        pdf_draw_slide(c, data, i + 1, total)
        c.showPage()
    c.save()
    print(f"✓ PDF saved:  {output_path}")


# ─────────────────────────────────────────────────────────────
# エントリーポイント
# ─────────────────────────────────────────────────────────────
if __name__ == "__main__":
    base = "/Users/maedaatsuya/projects/first-business"
    pptx_path = os.path.join(base, "pitch_deck.pptx")
    pdf_path  = os.path.join(base, "pitch_deck.pdf")
    build_pptx(pptx_path)
    build_pdf(pdf_path)
    print("\n完了！")
    print(f"  PPTX → {pptx_path}")
    print(f"  PDF  → {pdf_path}")
